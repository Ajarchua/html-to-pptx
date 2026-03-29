#!/usr/bin/env python3
"""
HTML → PPTX 元素映射转换器 v3
新增：
  - placeholder 机制：HTML 中 class="placeholder" 的 div 用 data-chart 属性描述图表，
    转换时用 python-pptx 图表 API 精确填入，坐标更可靠
  - 溢出检测：转换前扫描所有绝对定位元素，超出 1280×720 的给出警告

用法:
    python convert.py --input slides.html --output presentation.pptx
"""

import argparse
import json
import re
import sys
from pathlib import Path

from bs4 import BeautifulSoup
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import ChartData
from io import BytesIO

# ── 坐标系 ────────────────────────────────────────────────────────────────────
SLIDE_W_PX  = 1280
SLIDE_H_PX  = 720
SLIDE_W_EMU = 12192000
SLIDE_H_EMU = 6858000

def px_to_emu(px_val, axis='x'):
    if axis == 'x':
        return int(px_val * SLIDE_W_EMU / SLIDE_W_PX)
    return int(px_val * SLIDE_H_EMU / SLIDE_H_PX)

def px_to_pt(px_val):
    return px_val * 0.75


# ── CSS 解析工具 ──────────────────────────────────────────────────────────────
def parse_style(style_str):
    result = {}
    if not style_str:
        return result
    for part in style_str.split(';'):
        part = part.strip()
        if ':' in part:
            k, v = part.split(':', 1)
            result[k.strip().lower()] = v.strip()
    return result

def parse_px(val_str):
    if not val_str:
        return 0.0
    val_str = val_str.strip()
    if val_str.endswith('px'):
        try: return float(val_str[:-2])
        except ValueError: return 0.0
    try: return float(val_str)
    except ValueError: return 0.0

def parse_color(color_str):
    if not color_str or color_str in ('transparent', 'none', 'inherit'):
        return None
    color_str = color_str.strip()
    if color_str.startswith('#'):
        h = color_str[1:]
        if len(h) == 3: h = ''.join(c*2 for c in h)
        if len(h) == 6:
            try: return RGBColor(int(h[0:2],16), int(h[2:4],16), int(h[4:6],16))
            except ValueError: pass
    m = re.match(r'rgba?\s*\(\s*(\d+)\s*,\s*(\d+)\s*,\s*(\d+)', color_str)
    if m:
        return RGBColor(int(m.group(1)), int(m.group(2)), int(m.group(3)))
    named = {
        'white':(255,255,255),'black':(0,0,0),'red':(255,0,0),
        'green':(0,128,0),'blue':(0,0,255),'yellow':(255,255,0),
        'orange':(255,165,0),'purple':(128,0,128),
        'gray':(128,128,128),'grey':(128,128,128),
        'lightgray':(211,211,211),'darkgray':(169,169,169),
    }
    if color_str.lower() in named:
        return RGBColor(*named[color_str.lower()])
    return None

def parse_font_size(style):
    return px_to_pt(parse_px(style.get('font-size', '16px')))

def parse_font_weight(style):
    return style.get('font-weight', 'normal') in ('bold','700','800','900')

def parse_font_italic(style):
    return style.get('font-style', 'normal') == 'italic'

def parse_text_align(style):
    mapping = {'left':PP_ALIGN.LEFT,'center':PP_ALIGN.CENTER,
               'right':PP_ALIGN.RIGHT,'justify':PP_ALIGN.JUSTIFY}
    return mapping.get(style.get('text-align','left'), PP_ALIGN.LEFT)

def parse_background_color(style):
    bg = style.get('background-color') or style.get('background', '')
    if 'gradient' in bg.lower():
        m = re.search(r'#[0-9a-fA-F]{3,6}|rgb\([^)]+\)', bg)
        if m: return parse_color(m.group(0))
        return None
    return parse_color(bg)

def parse_geometry(style):
    left = px_to_emu(parse_px(style.get('left','0')), 'x')
    top  = px_to_emu(parse_px(style.get('top', '0')), 'y')
    w    = px_to_emu(parse_px(style.get('width','100')), 'x')
    h    = px_to_emu(parse_px(style.get('height','50')), 'y')
    return left, top, w, h


# ── 溢出检测 ──────────────────────────────────────────────────────────────────
def check_overflow(section, slide_index):
    """扫描幻灯片中所有绝对定位元素，超出边界则打印警告"""
    warnings = []
    for el in section.find_all(True, recursive=True):
        style = parse_style(el.get('style',''))
        if 'absolute' not in style.get('position',''):
            continue
        left   = parse_px(style.get('left',  '0'))
        top    = parse_px(style.get('top',   '0'))
        width  = parse_px(style.get('width', '0'))
        height = parse_px(style.get('height','0'))
        right  = left + width
        bottom = top  + height

        issues = []
        if right  > SLIDE_W_PX + 10: issues.append(f"右边超出 {right - SLIDE_W_PX:.0f}px")
        if bottom > SLIDE_H_PX + 10: issues.append(f"下边超出 {bottom - SLIDE_H_PX:.0f}px")
        if left   < -10:             issues.append(f"左边超出 {-left:.0f}px")
        if top    < -10:             issues.append(f"上边超出 {-top:.0f}px")

        if issues:
            tag   = el.name
            label = el.get('data-type') or ' '.join(el.get('class',[])) or tag
            text_preview = el.get_text(strip=True)[:30]
            warnings.append(
                f"  [溢出] 第{slide_index}张 <{tag} {label}> "
                f"({text_preview!r}): {', '.join(issues)}"
            )
    return warnings


# ── 元素处理器 ────────────────────────────────────────────────────────────────
def add_text_box(slide, el, style):
    left, top, w, h = parse_geometry(style)
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = parse_text_align(style)
    run = p.add_run()
    run.text = el.get_text()
    font = run.font
    font.size   = Pt(parse_font_size(style))
    font.bold   = parse_font_weight(style)
    font.italic = parse_font_italic(style)
    color = parse_color(style.get('color','#000000'))
    if color: font.color.rgb = color
    ff = style.get('font-family','Calibri').split(',')[0].strip().strip("'\"")
    font.name = ff


def add_shape(slide, el, style, shape_type):
    left, top, w, h = parse_geometry(style)
    type_map = {'rect':1, 'ellipse':9, 'arrow':13, 'right-arrow':13}
    mso_type = type_map.get(shape_type, 1)
    shape = slide.shapes.add_shape(mso_type, left, top, w, h)

    fill_color = parse_background_color(style)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()

    border_w = parse_px(style.get('border-width','0px'))
    border_c = parse_color(style.get('border-color',''))
    if border_w > 0 and border_c:
        shape.line.color.rgb = border_c
        shape.line.width = Pt(border_w * 0.75)
    else:
        shape.line.fill.background()

    text = el.get_text(strip=True)
    if text:
        tf = shape.text_frame
        tf.text = text
        for para in tf.paragraphs:
            para.alignment = parse_text_align(style)
            for run in para.runs:
                run.font.size = Pt(parse_font_size(style))
                c = parse_color(style.get('color','#000000'))
                if c: run.font.color.rgb = c


def add_image(slide, el, style):
    import requests
    left, top, w, h = parse_geometry(style)
    src = el.get('src','')
    if not src:
        return
    try:
        if src.startswith('data:image'):
            import base64
            _, data = src.split(',', 1)
            img_stream = BytesIO(base64.b64decode(data))
        elif src.startswith('http://') or src.startswith('https://'):
            resp = requests.get(src, timeout=10)
            img_stream = BytesIO(resp.content)
        else:
            img_stream = open(src, 'rb')
        slide.shapes.add_picture(img_stream, left, top, w, h)
    except Exception as e:
        print(f"  [警告] 图片加载失败: {src[:60]} ({e})", file=sys.stderr)


def add_table(slide, el, style):
    left, top, w, h = parse_geometry(style)
    rows_el  = el.find_all('tr')
    if not rows_el: return
    num_rows = len(rows_el)
    num_cols = max(len(r.find_all(['td','th'])) for r in rows_el)
    if not num_rows or not num_cols: return

    table = slide.shapes.add_table(num_rows, num_cols, left, top, w, h).table
    for ri, row_el in enumerate(rows_el):
        cells    = row_el.find_all(['td','th'])
        row_style = parse_style(row_el.get('style',''))
        row_bg   = parse_background_color(row_style)
        for ci, cell_el in enumerate(cells):
            if ci >= num_cols: break
            cell = table.cell(ri, ci)
            cell.text = cell_el.get_text(strip=True)
            cell_style  = parse_style(cell_el.get('style',''))
            merged      = {**row_style, **cell_style}
            bg = parse_background_color(merged) or row_bg
            if bg:
                cell.fill.solid()
                cell.fill.fore_color.rgb = bg
            for para in cell.text_frame.paragraphs:
                para.alignment = parse_text_align(merged)
                for run in para.runs:
                    run.font.size = Pt(parse_font_size(merged))
                    run.font.bold = parse_font_weight(merged)
                    c = parse_color(merged.get('color','#000000'))
                    if c: run.font.color.rgb = c
                    ff = merged.get('font-family','Calibri').split(',')[0].strip().strip("'\"")
                    run.font.name = ff


def add_chart_from_placeholder(slide, el, style):
    """
    Placeholder 机制：从 data-chart 属性读取图表定义，用 python-pptx 图表 API 填入。
    HTML 写法：
      <div class="placeholder" data-chart-type="bar"
           style="position:absolute; left:Xpx; top:Ypx; width:Wpx; height:Hpx;"
           data-chart='{"title":"...", "categories":[...], "series":[...]}'>
      </div>
    """
    left, top, w, h = parse_geometry(style)

    chart_type_str = el.get('data-chart-type', el.get('data-type', 'bar'))
    data_str       = el.get('data-chart', '{}')

    try:
        chart_def = json.loads(data_str)
    except json.JSONDecodeError:
        print(f"  [警告] placeholder 图表数据解析失败: {data_str[:60]}", file=sys.stderr)
        return

    chart_type_map = {
        'bar':            XL_CHART_TYPE.COLUMN_CLUSTERED,
        'column':         XL_CHART_TYPE.COLUMN_CLUSTERED,
        'line':           XL_CHART_TYPE.LINE,
        'pie':            XL_CHART_TYPE.PIE,
        'area':           XL_CHART_TYPE.AREA,
        'bar-horizontal': XL_CHART_TYPE.BAR_CLUSTERED,
        'scatter':        XL_CHART_TYPE.XY_SCATTER,
    }
    xl_type    = chart_type_map.get(chart_type_str, XL_CHART_TYPE.COLUMN_CLUSTERED)
    categories = chart_def.get('categories', [])
    series_list= chart_def.get('series', [])
    if not categories or not series_list:
        return

    cd = ChartData()
    cd.categories = categories
    for s in series_list:
        cd.add_series(s.get('name',''), s.get('values',[]))

    chart = slide.shapes.add_chart(xl_type, left, top, w, h, cd).chart

    title_text = chart_def.get('title','')
    if title_text:
        chart.has_title = True
        chart.chart_title.text_frame.text = title_text

    # 图例
    show_legend = chart_def.get('show_legend', len(series_list) > 1)
    chart.has_legend = show_legend


def add_chart_inline(slide, el, style):
    """旧的 data-chart 内嵌方式，保持向后兼容"""
    add_chart_from_placeholder(slide, el, style)


# ── 幻灯片背景 ────────────────────────────────────────────────────────────────
def set_slide_background(slide, style):
    bg_color = parse_background_color(style)
    if bg_color:
        fill = slide.background.fill
        fill.solid()
        fill.fore_color.rgb = bg_color


# ── 容器包裹写法处理 ──────────────────────────────────────────────────────────────
def add_shape_with_children(slide, shape_el, shape_style):
    """外层 shape 负责背景+绝对定位，内层 text-box 用 padding 描述偏移。"""
    left, top, w, h = parse_geometry(shape_style)
    shape_type = shape_el.get("data-type", "rect")
    type_map = {"rect": 1, "ellipse": 9, "arrow": 13}
    mso_type = type_map.get(shape_type, 1)

    shape = slide.shapes.add_shape(mso_type, left, top, w, h)
    fill_color = parse_background_color(shape_style)
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    shape.line.fill.background()

    cursor_top = top  # 当前行起点（EMU）

    for child in shape_el.find_all(["div", "p", "span"], recursive=True):
        child_cls   = " ".join(child.get("class", []))
        child_style = parse_style(child.get("style", ""))
        child_dtype = child.get("data-type", "")
        if child_dtype != "text" and "text-box" not in child_cls:
            continue
        text = child.get_text(strip=True)
        if not text:
            continue

        pt_px        = parse_px(child_style.get("padding-top", child_style.get("padding", "0")))
        pl_px        = parse_px(child_style.get("padding-left", child_style.get("padding", "16")))
        font_size_px = parse_px(child_style.get("font-size", "16px"))
        try:
            lh = float(str(child_style.get("line-height", "1.3")).rstrip("px"))
        except Exception:
            lh = 1.3

        txt_left = left + px_to_emu(pl_px, "x")
        txt_top  = cursor_top + px_to_emu(pt_px, "y")
        txt_w    = max(px_to_emu(10, "x"), w - px_to_emu(pl_px * 2, "x"))
        txt_h    = max(px_to_emu(font_size_px * lh * 1.4, "y"), px_to_emu(font_size_px * 2, "y"))

        if txt_top + txt_h > top + h + px_to_emu(20, "y"):
            break

        txBox = slide.shapes.add_textbox(txt_left, txt_top, txt_w, txt_h)
        tf    = txBox.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = parse_text_align(child_style)
        run = p.add_run()
        run.text    = text
        font        = run.font
        font.size   = Pt(px_to_pt(font_size_px))
        font.bold   = parse_font_weight(child_style)
        font.italic = parse_font_italic(child_style)
        color = parse_color(child_style.get("color", "#ffffff"))
        if color:
            font.color.rgb = color
        ff = child_style.get("font-family", "Calibri").split(",")[0].strip().strip("'\"")
        font.name = ff
        cursor_top = txt_top + txt_h


def _has_child_textboxes(el):
    """判断 shape 内是否有用 padding 定位的子文字框（容器包裹写法）"""
    for child in el.find_all(["div", "p", "span"], recursive=True):
        child_cls = " ".join(child.get("class", []))
        if "text-box" in child_cls or child.get("data-type") == "text":
            child_style = parse_style(child.get("style", ""))
            if "absolute" not in child_style.get("position", ""):
                return True
    return False


# ── 主转换逻辑 ────────────────────────────────────────────────────────────────
def process_slide(prs, section):
    blank_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank_layout)

    slide_style = parse_style(section.get("style", ""))
    set_slide_background(slide, slide_style)

    skip_els = set()

    for el in section.find_all(True, recursive=True):
        if id(el) in skip_els:
            continue

        el_style = parse_style(el.get("style", ""))
        dtype    = el.get("data-type", "")
        cls      = " ".join(el.get("class", []))

        if "absolute" not in el_style.get("position", ""):
            if el.name != "table" or "pptx-table" not in cls:
                continue

        # placeholder 图表
        if "placeholder" in cls and el.get("data-chart"):
            add_chart_from_placeholder(slide, el, el_style)

        # 文字框
        elif dtype == "text" or "text-box" in cls:
            add_text_box(slide, el, el_style)

        # 形状：检测是否为容器包裹写法
        elif dtype in ("rect", "ellipse", "arrow", "right-arrow") or "shape" in cls:
            if _has_child_textboxes(el):
                add_shape_with_children(slide, el, el_style)
                for child in el.find_all(True, recursive=True):
                    skip_els.add(id(child))
            else:
                add_shape(slide, el, el_style, dtype or "rect")

        # 图片
        elif el.name == "img":
            add_image(slide, el, el_style)

        # 表格
        elif el.name == "table" and "pptx-table" in cls:
            add_table(slide, el, el_style)

        # 旧式内嵌图表（向后兼容）
        elif "chart-box" in cls or dtype in ("bar", "line", "pie", "area", "column"):
            add_chart_inline(slide, el, el_style)

    return slide


def convert(html_path: str, output_path: str):
    print(f"读取 HTML: {html_path}")
    html  = Path(html_path).read_text(encoding='utf-8')
    soup  = BeautifulSoup(html, 'lxml')
    slides = soup.find_all('section', class_='slide')

    if not slides:
        print("错误：未找到 <section class='slide'> 元素", file=sys.stderr)
        sys.exit(1)

    print(f"发现 {len(slides)} 张幻灯片")

    # ── 溢出检测（转换前扫描全部幻灯片）──
    all_warnings = []
    for i, section in enumerate(slides, 1):
        all_warnings.extend(check_overflow(section, i))

    if all_warnings:
        print(f"\n⚠️  发现 {len(all_warnings)} 个元素超出幻灯片边界：")
        for w in all_warnings:
            print(w)
        print("  建议修正后重新生成，超出部分在 PPT 中可能被裁切。\n")
    else:
        print("✅ 溢出检测通过，所有元素在边界内")

    prs = Presentation()
    prs.slide_width  = Emu(SLIDE_W_EMU)
    prs.slide_height = Emu(SLIDE_H_EMU)

    for i, section in enumerate(slides, 1):
        print(f"  处理第 {i} 张...")
        process_slide(prs, section)

    prs.save(output_path)
    print(f"\n✅ 转换完成：{output_path}")


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='HTML → PPTX 转换器 v3')
    parser.add_argument('--input',  required=True, help='输入 HTML 文件路径')
    parser.add_argument('--output', required=True, help='输出 PPTX 文件路径')
    args = parser.parse_args()
    convert(args.input, args.output)
