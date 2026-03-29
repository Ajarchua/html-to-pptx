#!/usr/bin/env python3
"""
PPTX 模板风格提取器
分析用户上传的 .pptx 文件第一张幻灯片，提取配色、字体、布局规律，
输出 style_profile.json 供 AI 生成 HTML 时参考。

用法:
    python extract_style.py --input template.pptx --output style_profile.json
"""

import argparse
import json
import sys
from collections import Counter
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE_TYPE


def rgb_to_hex(rgb):
    return f"#{rgb[0]:02x}{rgb[1]:02x}{rgb[2]:02x}"


def get_shape_fill_color(shape):
    try:
        fill = shape.fill
        if fill.type is None:
            return None
        if hasattr(fill, 'fore_color') and fill.fore_color:
            try:
                return rgb_to_hex(fill.fore_color.rgb)
            except Exception:
                pass
    except Exception:
        pass
    return None


def get_text_color(run):
    try:
        if run.font.color and run.font.color.rgb:
            return rgb_to_hex(run.font.color.rgb)
    except Exception:
        pass
    return None


def get_background_color(slide, prs):
    try:
        bg = slide.background
        fill = bg.fill
        if fill.fore_color:
            return rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    try:
        bg = slide.slide_layout.background
        fill = bg.fill
        if fill.fore_color:
            return rgb_to_hex(fill.fore_color.rgb)
    except Exception:
        pass
    return "#ffffff"


def analyze_slide(slide, prs):
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    colors_found, text_colors, title_fonts, body_fonts, shapes_info = [], [], [], [], []

    for shape in slide.shapes:
        left_pct   = shape.left   / slide_w if slide_w else 0
        top_pct    = shape.top    / slide_h if slide_h else 0
        width_pct  = shape.width  / slide_w if slide_w else 0
        height_pct = shape.height / slide_h if slide_h else 0

        fill_color = get_shape_fill_color(shape)
        if fill_color and fill_color != "#ffffff":
            colors_found.append(fill_color)
            shapes_info.append({
                "type": _shape_type_str(shape),
                "position": _position_label(left_pct, top_pct, width_pct, height_pct),
                "color": fill_color,
                "left_pct":   round(left_pct, 3),
                "top_pct":    round(top_pct, 3),
                "width_pct":  round(width_pct, 3),
                "height_pct": round(height_pct, 3),
            })

        if shape.has_text_frame:
            for para in shape.text_frame.paragraphs:
                for run in para.runs:
                    if not run.text.strip():
                        continue
                    tc = get_text_color(run)
                    if tc:
                        text_colors.append(tc)
                    font_name = run.font.name or "Calibri"
                    font_size = run.font.size
                    font_size_pt = round(font_size / 12700) if font_size else None
                    is_bold = run.font.bold or False
                    entry = {"name": font_name, "size_pt": font_size_pt, "bold": is_bold}
                    if font_size_pt and font_size_pt >= 24:
                        title_fonts.append(entry)
                    elif font_size_pt:
                        body_fonts.append(entry)

    return colors_found, text_colors, title_fonts, body_fonts, shapes_info


def _shape_type_str(shape):
    try:
        st = shape.shape_type
        if st == MSO_SHAPE_TYPE.AUTO_SHAPE: return "rect"
        elif st == MSO_SHAPE_TYPE.PICTURE:  return "image"
        elif st == MSO_SHAPE_TYPE.TEXT_BOX: return "text"
        else: return "shape"
    except Exception:
        return "shape"


def _position_label(left, top, width, height):
    if left < 0.1 and width < 0.4:          return "left"
    elif left > 0.6 and width < 0.4:        return "right"
    elif top < 0.15 and height < 0.3:       return "top"
    elif top > 0.7:                          return "bottom"
    elif left < 0.1 and width > 0.9:        return "full-width"
    else:                                    return "center"


def pick_dominant(lst):
    if not lst: return None
    return Counter(lst).most_common(1)[0][0]


def pick_best_font(font_list, fallback_name, fallback_size, fallback_bold):
    if not font_list:
        return {"name": fallback_name, "size_pt": fallback_size, "bold": fallback_bold}
    with_size = [f for f in font_list if f.get("size_pt")]
    return with_size[0] if with_size else font_list[0]


def infer_layout(shapes_info):
    layout = {
        "has_left_sidebar": False, "sidebar_width_pct": 0,
        "has_top_bar": False,      "top_bar_height_pct": 0,
        "decorative_shapes": [],
    }
    for s in shapes_info:
        lp, tp, wp, hp = s["left_pct"], s["top_pct"], s["width_pct"], s["height_pct"]
        if lp < 0.05 and 0.1 < wp < 0.5 and hp > 0.6:
            layout["has_left_sidebar"]   = True
            layout["sidebar_width_pct"]  = round(wp * 100)
            layout["decorative_shapes"].append({"type": s["type"], "position": "left",  "color": s["color"]})
        elif tp < 0.05 and hp < 0.25 and wp > 0.7:
            layout["has_top_bar"]        = True
            layout["top_bar_height_pct"] = round(hp * 100)
            layout["decorative_shapes"].append({"type": s["type"], "position": "top",   "color": s["color"]})
        elif wp * hp < 0.15:
            layout["decorative_shapes"].append({"type": s["type"], "position": s["position"], "color": s["color"]})
    return layout


def _is_dark(hex_color):
    try:
        h = hex_color.lstrip('#')
        r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        return (0.299*r + 0.587*g + 0.114*b) < 128
    except Exception:
        return False


def _color_dist(c1, c2):
    try:
        h1,h2 = c1.lstrip('#'), c2.lstrip('#')
        r1,g1,b1 = int(h1[0:2],16), int(h1[2:4],16), int(h1[4:6],16)
        r2,g2,b2 = int(h2[0:2],16), int(h2[2:4],16), int(h2[4:6],16)
        return ((r1-r2)**2+(g1-g2)**2+(b1-b2)**2)**0.5
    except Exception:
        return 0


def _lighten_or_darken(hex_color):
    try:
        h = hex_color.lstrip('#')
        r,g,b = int(h[0:2],16), int(h[2:4],16), int(h[4:6],16)
        if _is_dark(hex_color):
            r,g,b = min(255,r+80), min(255,g+80), min(255,b+80)
        else:
            r,g,b = max(0,r-80), max(0,g-80), max(0,b-80)
        return f"#{r:02x}{g:02x}{b:02x}"
    except Exception:
        return "#888888"


def extract(pptx_path, output_path):
    print(f"读取模板: {pptx_path}")
    prs = Presentation(pptx_path)
    if not prs.slides:
        print("错误：模板中没有幻灯片", file=sys.stderr)
        sys.exit(1)

    slide = prs.slides[0]
    print(f"分析第一张幻灯片（共 {len(prs.slides)} 张）...")

    bg_color = get_background_color(slide, prs)
    colors_found, text_colors, title_fonts, body_fonts, shapes_info = analyze_slide(slide, prs)

    non_text  = [c for c in colors_found if c not in text_colors]
    primary   = pick_dominant(non_text) or pick_dominant(colors_found) or "#1a1a2e"
    unique_nt = list(set(non_text))
    accent    = (max(unique_nt, key=lambda c: _color_dist(c, primary)) if len(unique_nt) >= 2 else None) or "#e94560"
    text_main = pick_dominant(text_colors) or ("#ffffff" if _is_dark(bg_color) else "#1a1a1a")
    sec_list  = [c for c in set(text_colors) if c != text_main]
    text_secondary = sec_list[0] if sec_list else _lighten_or_darken(text_main)

    title_font = pick_best_font(title_fonts, "Microsoft YaHei", 40, True)
    body_font  = pick_best_font(body_fonts,  "Calibri",         16, False)
    layout     = infer_layout(shapes_info)

    profile = {
        "source": Path(pptx_path).name,
        "slide_size": {"width_emu": prs.slide_width, "height_emu": prs.slide_height},
        "colors": {
            "background":     bg_color,
            "primary":        primary,
            "accent":         accent,
            "text_main":      text_main,
            "text_secondary": text_secondary,
        },
        "fonts": {"title": title_font, "body": body_font},
        "layout": layout,
        "_raw_colors_found": list(set(colors_found))[:20],
    }

    Path(output_path).write_text(json.dumps(profile, ensure_ascii=False, indent=2))
    print(f"✅ 风格提取完成：{output_path}")
    print(f"   背景色: {bg_color}  主色: {primary}  强调色: {accent}")
    print(f"   标题字体: {title_font['name']} {title_font.get('size_pt','?')}pt")
    print(f"   正文字体: {body_font['name']}  {body_font.get('size_pt','?')}pt")
    if layout["has_left_sidebar"]:
        print(f"   检测到左侧色块，宽度约 {layout['sidebar_width_pct']}%")
    if layout["has_top_bar"]:
        print(f"   检测到顶部色条，高度约 {layout['top_bar_height_pct']}%")


if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='PPTX 模板风格提取器')
    parser.add_argument('--input',  required=True, help='输入 .pptx 模板路径')
    parser.add_argument('--output', required=True, help='输出 style_profile.json 路径')
    args = parser.parse_args()
    extract(args.input, args.output)
